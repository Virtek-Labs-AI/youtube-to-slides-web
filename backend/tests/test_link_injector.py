import io

from pptx import Presentation

from app.services.link_injector import inject_references


def _pptx_with_bullets() -> bytes:
    """Build a PPTX with slides that have title + bullet text."""
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title + Content

    # Slide 0: title slide
    slide0 = prs.slides.add_slide(layout)
    slide0.shapes.title.text = "My Video"

    # Slide 1: content slide with 2 bullets
    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Topic A"
    tf = slide1.placeholders[1].text_frame
    tf.text = "First point"
    para = tf.add_paragraph()
    para.text = "Second point"

    # Slide 2: content slide with 1 bullet
    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Topic B"
    tf = slide2.placeholders[1].text_frame
    tf.text = "Linked point"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


REFERENCES = [
    {"slide": 1, "bullet": 0, "url": "https://youtu.be/abc123?t=60"},
    {"slide": 1, "bullet": 1, "url": "https://youtu.be/abc123?t=120"},
    {"slide": 2, "bullet": 0, "url": "https://youtu.be/abc123?t=300"},
]


class TestInjectReferences:
    def test_no_references_returns_unchanged(self) -> None:
        original = _pptx_with_bullets()
        result = inject_references(original, [])
        assert result == original

    def test_returns_bytes(self) -> None:
        original = _pptx_with_bullets()
        result = inject_references(original, REFERENCES)
        assert isinstance(result, bytes)
        assert len(result) > 0

    def test_injects_hyperlinks_into_bullets(self) -> None:
        original = _pptx_with_bullets()
        result = inject_references(original, REFERENCES)
        prs = Presentation(io.BytesIO(result))

        # Slide 1 (index 1) should have hyperlinks on both bullets
        slide = prs.slides[1]
        hyperlinks = _get_hyperlinks(slide)
        assert "https://youtu.be/abc123?t=60" in hyperlinks
        assert "https://youtu.be/abc123?t=120" in hyperlinks

    def test_injects_on_correct_slide(self) -> None:
        original = _pptx_with_bullets()
        result = inject_references(original, REFERENCES)
        prs = Presentation(io.BytesIO(result))

        # Slide 2 (index 2) should have hyperlink
        slide = prs.slides[2]
        hyperlinks = _get_hyperlinks(slide)
        assert "https://youtu.be/abc123?t=300" in hyperlinks

    def test_slide_count_unchanged(self) -> None:
        original = _pptx_with_bullets()
        result = inject_references(original, REFERENCES)
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 3  # no extra slides added

    def test_skips_out_of_range_slide(self) -> None:
        original = _pptx_with_bullets()
        refs = [{"slide": 99, "bullet": 0, "url": "https://youtu.be/abc?t=1"}]
        result = inject_references(original, refs)
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 3

    def test_skips_out_of_range_bullet(self) -> None:
        original = _pptx_with_bullets()
        refs = [{"slide": 1, "bullet": 99, "url": "https://youtu.be/abc?t=1"}]
        result = inject_references(original, refs)
        # Should not raise, just skip
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 3

    def test_skips_incomplete_references(self) -> None:
        original = _pptx_with_bullets()
        refs = [
            {"slide": 1, "bullet": 0},  # missing url
            {"slide": 1, "url": "https://youtu.be/abc?t=1"},  # missing bullet
            {"bullet": 0, "url": "https://youtu.be/abc?t=1"},  # missing slide
        ]
        result = inject_references(original, refs)
        prs = Presentation(io.BytesIO(result))
        # No hyperlinks should be added
        slide = prs.slides[1]
        hyperlinks = _get_hyperlinks(slide)
        assert len(hyperlinks) == 0


def _get_hyperlinks(slide) -> list[str]:
    """Extract all hyperlink addresses from non-title shapes on a slide."""
    hyperlinks = []
    title_shape_id = None
    if slide.shapes.title:
        title_shape_id = slide.shapes.title.shape_id

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if title_shape_id and shape.shape_id == title_shape_id:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.hyperlink and run.hyperlink.address:
                    hyperlinks.append(run.hyperlink.address)
    return hyperlinks
