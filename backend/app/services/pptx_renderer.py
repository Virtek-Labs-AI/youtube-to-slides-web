import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.core.config import settings


def render_pptx(slides_data: dict, filename: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = slides_data.get("slides", [])

    for slide_data in slides:
        slide_type = slide_data.get("type", "content")
        if slide_type == "title":
            _add_title_slide(prs, slide_data)
        elif slide_type in ("content", "overview", "takeaway"):
            _add_content_slide(prs, slide_data)

    storage_dir = settings.storage_path
    os.makedirs(storage_dir, exist_ok=True)
    filepath = os.path.join(storage_dir, filename)
    prs.save(filepath)
    return filepath


def _add_title_slide(prs: Presentation, data: dict) -> None:
    layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(layout)

    if slide.placeholders[0]:
        slide.placeholders[0].text = data.get("title", "")
    if len(slide.placeholders) > 1 and slide.placeholders[1]:
        slide.placeholders[1].text = data.get("subtitle", "")


def _add_content_slide(prs: Presentation, data: dict) -> None:
    layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(layout)

    # Set title
    if slide.placeholders[0]:
        slide.placeholders[0].text = data.get("title", "")

    # Build bullet content in the body placeholder
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    bullets = data.get("bullets", [])
    for i, bullet in enumerate(bullets):
        text = bullet.get("text", "") if isinstance(bullet, dict) else str(bullet)
        url = bullet.get("url", "") if isinstance(bullet, dict) else ""

        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        # Add bullet text with hyperlink if URL exists
        run = para.add_run()
        run.text = text
        run.font.size = Pt(18)

        if url:
            run.hyperlink.address = url

            # Add timestamp label in gray
            ts_run = para.add_run()
            ts_label = _extract_timestamp_label(url)
            ts_run.text = f"  {ts_label}"
            ts_run.font.size = Pt(14)
            ts_run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


def _extract_timestamp_label(url: str) -> str:
    if "?t=" in url:
        try:
            seconds = int(url.split("?t=")[1])
            minutes, secs = divmod(seconds, 60)
            return f"({minutes}:{secs:02d})"
        except (ValueError, IndexError):
            pass
    return ""
