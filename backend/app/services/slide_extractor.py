"""Extract slide text from a Presenton-generated PPTX.

Reads the PPTX and returns a structured list of slides with their titles
and bullet text, ready for reference matching.
"""

import io

from pptx import Presentation


def extract_slides(pptx_bytes: bytes) -> list[dict]:
    """Extract slide titles and bullet text from PPTX bytes.

    Returns a list of dicts: [{"slide": 0, "title": "...", "bullets": [{"index": 0, "text": "..."}]}]
    """
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides = []

    for i, slide in enumerate(prs.slides):
        title = ""
        bullets = []

        title_shape_id = None
        if slide.shapes.title:
            title = slide.shapes.title.text_frame.text.strip()
            title_shape_id = slide.shapes.title.shape_id

        bullet_idx = 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if title_shape_id and shape.shape_id == title_shape_id:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    bullets.append({"index": bullet_idx, "text": text})
                    bullet_idx += 1

        if not title and bullets:
            title = bullets.pop(0)["text"]

        slides.append({"slide": i, "title": title, "bullets": bullets})

    return slides
