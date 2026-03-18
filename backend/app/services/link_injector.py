"""Inject YouTube reference links into a Presenton-generated PPTX.

Takes a list of references (slide index, bullet index, url) and adds
hyperlinks directly to the matching bullet text runs in the PPTX.
"""

import io

from pptx import Presentation


def inject_references(pptx_bytes: bytes, references: list[dict]) -> bytes:
    """Add hyperlinks to bullet text in a PPTX based on reference mappings.

    Each reference dict has: {"slide": int, "bullet": int, "url": str}
    Hyperlinks are added directly to the first run of the matching paragraph.

    Returns the modified PPTX as bytes.
    """
    if not references:
        return pptx_bytes

    prs = Presentation(io.BytesIO(pptx_bytes))
    slides = list(prs.slides)

    for ref in references:
        slide_idx = ref.get("slide")
        bullet_idx = ref.get("bullet")
        url = ref.get("url")

        if slide_idx is None or bullet_idx is None or not url:
            continue
        if slide_idx >= len(slides):
            continue

        slide = slides[slide_idx]

        # Collect non-title text paragraphs
        paragraphs = []
        title_shape_id = None
        if slide.shapes.title:
            title_shape_id = slide.shapes.title.shape_id

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if title_shape_id and shape.shape_id == title_shape_id:
                continue
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    paragraphs.append(para)

        if bullet_idx >= len(paragraphs):
            continue

        para = paragraphs[bullet_idx]

        # Add hyperlink to the first run, or create one if none exist
        if para.runs:
            para.runs[0].hyperlink.address = url
        else:
            run = para.add_run()
            run.text = para.text
            run.hyperlink.address = url

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
